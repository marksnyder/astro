"""Document loading and chunking for supported file types."""

from pathlib import Path

from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc",
    ".pptx", ".xlsx", ".xls",
    ".txt", ".md", ".csv",
}


# ── Custom loaders (avoids the heavy 'unstructured' dependency) ──────────


def _load_pptx(file_path: str) -> list[Document]:
    prs = Presentation(file_path)
    docs: list[Document] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        if lines:
            docs.append(
                Document(
                    page_content="\n".join(lines),
                    metadata={"source": file_path, "slide": slide_num},
                )
            )
    return docs


def _load_xlsx(file_path: str) -> list[Document]:
    wb = load_workbook(file_path, data_only=True)
    docs: list[Document] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            docs.append(
                Document(
                    page_content="\n".join(rows),
                    metadata={"source": file_path, "sheet": sheet_name},
                )
            )
    return docs


# ── Public API ───────────────────────────────────────────────────────────


def load_document(file_path: str) -> list[Document]:
    """Load a single file and return a list of Documents."""
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return PyPDFLoader(file_path).load()
    if ext in (".docx", ".doc"):
        return Docx2txtLoader(file_path).load()
    if ext == ".pptx":
        return _load_pptx(file_path)
    if ext in (".xlsx", ".xls"):
        return _load_xlsx(file_path)
    if ext in (".txt", ".md", ".csv"):
        return TextLoader(file_path, encoding="utf-8").load()

    print(f"  Skipping unsupported file: {file_path}")
    return []


def load_documents(path: str) -> list[Document]:
    """Load all supported documents from a file or directory."""
    p = Path(path)
    if p.is_file():
        files = [p]
    elif p.is_dir():
        files = [f for ext in SUPPORTED_EXTENSIONS for f in p.rglob(f"*{ext}")]
    else:
        raise FileNotFoundError(f"Path not found: {path}")

    documents: list[Document] = []
    for f in sorted(set(files)):
        print(f"  Loading: {f.name}")
        try:
            documents.extend(load_document(str(f)))
        except Exception as e:
            print(f"  Error loading {f.name}: {e}")

    return documents


def chunk_documents(
    documents: list[Document],
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[Document]:
    """Split documents into chunks for embedding."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return splitter.split_documents(documents)
